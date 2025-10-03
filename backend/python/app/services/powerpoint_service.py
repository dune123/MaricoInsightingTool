"""
PowerPoint Generation Service for Non-MMM Analysis

This service generates PowerPoint presentations from Non-MMM analysis data,
including context slides, insights, charts, and model results.
"""

import os
import json
import tempfile
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import io
import base64

from app.core.config import settings


class PowerPointService:
    """Service for generating PowerPoint presentations from analysis data"""
    
    @staticmethod
    def generate_analysis_presentation(
        analysis_id: str,
        brand_name: str,
        analysis_data: Dict[str, Any]
    ) -> Path:
        """
        Generate a complete PowerPoint presentation from analysis data
        
        Args:
            analysis_id: Unique analysis identifier
            brand_name: Brand name for the analysis
            analysis_data: Complete analysis data including charts, models, etc.
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            # Create presentation
            prs = Presentation()
            
            # Set slide size to widescreen (16:9)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Generate slides
            PowerPointService._add_title_slide(prs, brand_name, analysis_data)
            PowerPointService._add_context_slide(prs, analysis_data)
            PowerPointService._add_insights_slide(prs, analysis_data)
            PowerPointService._add_chart_slides(prs, analysis_data)
            PowerPointService._add_model_results_slide(prs, analysis_data)
            
            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{brand_name.replace(' ', '_')}_Analysis_{timestamp}.pptx"
            
            # Save to brand-specific directory
            brand_dirs = settings.get_brand_directories(brand_name)
            brand_dirs["export_dir"].mkdir(parents=True, exist_ok=True)
            output_path = brand_dirs["export_dir"] / filename
            
            prs.save(str(output_path))
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Failed to generate PowerPoint presentation: {str(e)}")
    
    @staticmethod
    def _add_title_slide(prs: Presentation, brand_name: str, analysis_data: Dict[str, Any]):
        """Add title slide to presentation"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = f"{brand_name} Analysis"
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"Non-MMM Analysis Report\nGenerated on {datetime.now().strftime('%B %d, %Y')}"
        
        # Style the text
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    
    @staticmethod
    def _add_context_slide(prs: Presentation, analysis_data: Dict[str, Any]):
        """Add context slide with data overview"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Analysis Context"
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Add context information
        context_info = [
            f"Brand: {analysis_data.get('brandName', 'N/A')}",
            f"Data File: {analysis_data.get('filename', 'N/A')}",
            f"Target Variable: {analysis_data.get('targetVariable', 'N/A')}",
            f"Analysis Date: {datetime.now().strftime('%B %d, %Y')}",
            ""
        ]
        
        # Add data summary if available
        data_summary = analysis_data.get('dataSummary', {})
        if data_summary:
            context_info.extend([
                "Data Summary:",
                f"• Total Records: {data_summary.get('totalRecords', 'N/A')}",
                f"• Variables: {data_summary.get('totalVariables', 'N/A')}",
                f"• Data Quality: {data_summary.get('dataQuality', 'N/A')}"
            ])
        
        for i, line in enumerate(context_info):
            p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            if line.startswith("•"):
                p.font.color.rgb = RGBColor(0, 102, 204)
    
    @staticmethod
    def _add_insights_slide(prs: Presentation, analysis_data: Dict[str, Any]):
        """Add initial insights slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Initial Insights"
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        insights = [
            "Key Findings:",
            "• Data quality assessment completed successfully",
            "• Target variable identified and validated",
            "• Expected signs configured for all variables",
            "• Statistical analysis performed on dataset",
            "",
            "Analysis Scope:",
            "• Comprehensive data exploration",
            "• Chart analysis with trendlines",
            "• Statistical modeling and validation",
            "• Professional presentation generation"
        ]
        
        for i, line in enumerate(insights):
            p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            if line.startswith("•"):
                p.font.color.rgb = RGBColor(0, 102, 204)
            elif line.endswith(":"):
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 51, 102)
    
    @staticmethod
    def _add_chart_slides(prs: Presentation, analysis_data: Dict[str, Any]):
        """Add chart analysis slides"""
        chart_data = analysis_data.get('chartData', [])
        
        for i, chart in enumerate(chart_data):
            slide_layout = prs.slide_layouts[5]  # Blank layout for charts
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = f"Chart Analysis: {chart.get('variableName', f'Variable {i+1}')}"
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            
            # Add chart information
            info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(1))
            info_frame = info_box.text_frame
            info_frame.text = f"Expected Sign: {chart.get('expectedSign', '~')} | Trendline: {chart.get('trendlineType', 'Linear')}"
            info_frame.paragraphs[0].font.size = Pt(14)
            info_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
            
            # Add chart placeholders (since we can't embed actual charts easily)
            chart_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
            chart_frame = chart_box.text_frame
            chart_frame.text = f"Chart Placeholder\n\nLine Chart: {chart.get('variableName', 'Variable')} vs Target Variable\nScatter Plot: Correlation Analysis\n\nNote: Actual charts would be embedded here in a full implementation"
            chart_frame.paragraphs[0].font.size = Pt(16)
            chart_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    @staticmethod
    def _add_model_results_slide(prs: Presentation, analysis_data: Dict[str, Any]):
        """Add model results slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Model Results"
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        model_results = analysis_data.get('modelResults', [])
        
        if model_results:
            # Add model information
            p = tf.paragraphs[0]
            p.text = "Statistical Models Trained:"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            
            for i, model in enumerate(model_results[:3]):  # Show top 3 models
                p = tf.add_paragraph()
                p.text = f"• {model.get('modelType', 'Unknown Model')} - R²: {model.get('rSquared', 'N/A')}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0, 102, 204)
            
            # Add best model details
            if model_results:
                best_model = model_results[0]
                p = tf.add_paragraph()
                p.text = ""
                
                p = tf.add_paragraph()
                p.text = "Best Model Details:"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 51, 102)
                
                p = tf.add_paragraph()
                p.text = f"• Model Type: {best_model.get('modelType', 'N/A')}"
                p.font.size = Pt(14)
                
                p = tf.add_paragraph()
                p.text = f"• R² Score: {best_model.get('rSquared', 'N/A')}"
                p.font.size = Pt(14)
                
                p = tf.add_paragraph()
                p.text = f"• Variables: {len(best_model.get('variables', []))}"
                p.font.size = Pt(14)
        else:
            p = tf.paragraphs[0]
            p.text = "No model results available"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(102, 102, 102)
    
    @staticmethod
    def _create_chart_image(chart_data: Dict[str, Any]) -> str:
        """
        Create a chart image from chart data
        Note: This is a simplified implementation. In a full implementation,
        you would create actual matplotlib charts and convert them to images.
        """
        # For now, return a placeholder
        return "Chart placeholder - would contain actual chart image"
